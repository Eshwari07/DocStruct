from __future__ import annotations

import datetime as _dt
from pathlib import Path
from typing import List

from pptx import Presentation  # type: ignore[import]
from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import]

from docstruct.core.schema import (
    DocumentTree,
    DocNode,
    PageRange,
    SourceFormat,
    ExtractionPath,
)
from docstruct.parsers.base import BaseParser


def _get_slide_title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.text:
        return slide.shapes.title.text.strip()
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text:
            return shape.text.strip()
    return ""


def _is_section_divider(slide) -> bool:
    title = _get_slide_title(slide)
    if not title or len(title.split()) > 8:
        return False
    # If there is other body text, treat as normal slide
    for shape in slide.shapes:
        if shape == slide.shapes.title:
            continue
        if hasattr(shape, "text") and shape.text and shape.text.strip():
            return False
    return True


class PptxParser(BaseParser):
    """
    Fast-path PPTX parser.

    - Uses slide titles as headings
    - Treats title-only slides as section dividers (parents)
    - Other slides become children of the current section
    """

    def parse(self) -> DocumentTree:
        pres = Presentation(str(self.file_path))

        root = DocNode(
            title=self.file_path.stem,
            text="",
            pages=PageRange(physical_start=1, physical_end=len(pres.slides)),
            confidence=1.0,
            images=[],
            children=[],
        )

        current_section = root

        for idx, slide in enumerate(pres.slides):
            page_num = idx + 1
            pages = PageRange(
                physical_start=page_num,
                physical_end=page_num,
                logical_start=None,
                logical_end=None,
            )

            title = _get_slide_title(slide) or f"Slide {page_num}"

            # Collect body text (non-title placeholders)
            body_parts: List[str] = []
            for shape in slide.shapes:
                if shape == slide.shapes.title:
                    continue
                if hasattr(shape, "text") and shape.text:
                    text = shape.text.strip()
                    if text:
                        body_parts.append(text)

            is_divider = _is_section_divider(slide)

            node = DocNode(
                title=title,
                text="\n\n".join(body_parts),
                pages=pages,
                confidence=1.0,
                images=[],
                children=[],
            )

            if is_divider:
                root.add_child(node)
                current_section = node
            else:
                current_section.add_child(node)

        tree = DocumentTree(
            source_file=str(self.file_path.name),
            source_format=SourceFormat.PPTX,
            total_pages=len(pres.slides),
            extracted_at=_dt.datetime.utcnow().isoformat() + "Z",
            extraction_path=ExtractionPath.FAST,
            nodes=[root],
        )
        tree.finalise()
        return tree

